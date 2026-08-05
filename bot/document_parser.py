import os
import logging
from typing import Optional

def parse_document(file_path: str, filename: str, max_chars: int = 50000) -> Optional[str]:
    """
    Parses a document based on its extension and extracts text.
    Returns the extracted text, truncated to max_chars if it exceeds it.
    Supported extensions: .pdf, .docx, .xlsx, .pptx, .txt
    """
    _, ext = os.path.splitext(filename.lower())
    text = ""
    
    try:
        if ext == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
                
        elif ext == '.pdf':
            import pypdf
            with open(file_path, 'rb') as f:
                reader = pypdf.PdfReader(f)
                for page in reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
                        
        elif ext == '.docx':
            import docx
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + "\t"
                    text += "\n"
                
        elif ext == '.xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"--- Sheet: {sheet_name} ---\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
                        
        elif ext == '.pptx':
            import pptx
            prs = pptx.Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                text += f"--- Slide {i + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                text += cell.text + "\t"
                            text += "\n"
        else:
            return None
            
        # Hard truncate to max_chars characters
        if len(text) > max_chars:
            text = text[:max_chars] + "\n\n[文件后续内容由于长度限制已被截断]"
            
        return text.strip()
        
    except Exception as e:
        logging.error(f"Error parsing document {filename}: {e}")
        return None
